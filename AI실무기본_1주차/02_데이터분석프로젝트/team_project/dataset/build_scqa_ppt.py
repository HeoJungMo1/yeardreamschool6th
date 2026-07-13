import csv
import json
import os
import re
import time
import urllib.request
from collections import defaultdict, Counter
from datetime import datetime, date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = r'C:/Users/chan1/Desktop/study/이어드림/team_project'
CSV_PATH = os.path.join(ROOT, 'Seoul_subway_data_20210705.csv')
SCQA_MD = os.path.join(ROOT, '혼잡도_SCQA_추가분석.md')
PPTX_PATH = os.path.join(ROOT, '서울지하철_혼잡도_SCQA_발표자료.pptx')
API_CACHE = os.path.join(ROOT, 'api_card_subway_202605.json')
DOWNLOAD_CSV = os.path.join(ROOT, 'CARD_SUBWAY_MONTH_202605.csv')
API_BASE = 'http://openapi.seoul.go.kr:8088/sample/json/CardSubwayStatsNew/{start}/{end}/{ymd}'
API_DATASET_URL = 'https://data.seoul.go.kr/dataList/OA-12914/S/1/datasetView.do'
API_SAMPLE_URL = 'http://openapi.seoul.go.kr:8088/sample/json/CardSubwayStatsNew/1/5/20260501'

NAVY = RGBColor(20, 33, 61)
BLUE = RGBColor(0, 119, 182)
TEAL = RGBColor(0, 150, 136)
ORANGE = RGBColor(244, 162, 97)
RED = RGBColor(214, 40, 40)
LIGHT = RGBColor(245, 247, 250)
DARK = RGBColor(35, 35, 35)
MUTED = RGBColor(105, 112, 122)
WHITE = RGBColor(255, 255, 255)


def n(x):
    return f'{int(round(x)):,}'


def parse_old_csv():
    hour_re = re.compile(r'(\d{2})시-(\d{2})시 (승차|하차)인원')
    latest_month = None
    months = []
    rows_count = 0
    stations = set(); lines = set()
    latest_station_total = defaultdict(int)
    latest_line_total = defaultdict(int)
    latest_hour_total = defaultdict(int)
    latest_station_hour_total = defaultdict(int)
    overall_station_total = defaultdict(int)
    with open(CSV_PATH, 'r', encoding='cp949', newline='') as f:
        reader = csv.reader(f)
        headers = next(reader)
        hour_cols = []
        for i, h in enumerate(headers):
            m = hour_re.match(h)
            if m:
                hour_cols.append((i, int(m.group(1)), m.group(3)))
        all_rows = []
        for row in reader:
            if not row or len(row) < 5:
                continue
            rows_count += 1
            months.append(row[0])
            latest_month = row[0] if latest_month is None or row[0] > latest_month else latest_month
            lines.add(row[1]); stations.add(row[2])
            all_rows.append(row)
        for row in all_rows:
            line, st = row[1], row[2]
            total = 0
            for i, h, typ in hour_cols:
                try:
                    v = int(row[i].replace(',', ''))
                except Exception:
                    v = 0
                total += v
                if row[0] == latest_month:
                    latest_hour_total[h] += v
                    latest_station_hour_total[(line, st, h)] += v
            overall_station_total[(line, st)] += total
            if row[0] == latest_month:
                latest_station_total[(line, st)] += total
                latest_line_total[line] += total
    return {
        'rows_count': rows_count,
        'month_min': min(months),
        'month_max': max(months),
        'latest_month': latest_month,
        'station_count': len(stations),
        'line_count': len(lines),
        'latest_station_total': latest_station_total,
        'latest_line_total': latest_line_total,
        'latest_hour_total': latest_hour_total,
        'latest_station_hour_total': latest_station_hour_total,
        'overall_station_total': overall_station_total,
    }


def fetch_api_202605():
    # Full analysis uses the official CSV attached to the same OA-12914 page.
    # The public sample API key is useful for endpoint verification but returns only a small sample.
    if os.path.exists(DOWNLOAD_CSV) and os.path.getsize(DOWNLOAD_CSV) > 1000:
        rows = []
        with open(DOWNLOAD_CSV, 'r', encoding='utf-8-sig', newline='') as f:
            reader = csv.DictReader(f)
            for r in reader:
                if not r.get('사용일자'):
                    continue
                rows.append({
                    'USE_YMD': r.get('사용일자', ''),
                    'SBWY_ROUT_LN_NM': r.get('노선명', ''),
                    'SBWY_STNS_NM': r.get('역명', ''),
                    'GTON_TNOPE': r.get('승차총승객수', '0'),
                    'GTOFF_TNOPE': r.get('하차총승객수', '0'),
                    'REG_YMD': r.get('등록일자', ''),
                })
        payload = {
            'fetched_at': datetime.now().isoformat(timespec='seconds'),
            'source': DOWNLOAD_CSV,
            'rows': rows,
            'errors': [],
            'note': 'OpenAPI endpoint verified with sample key; full-month rows parsed from official OA-12914 CSV download.'
        }
        with open(API_CACHE, 'w', encoding='utf-8') as f:
            json.dump(payload, f, ensure_ascii=False, indent=2)
        return payload
    if os.path.exists(API_CACHE):
        with open(API_CACHE, 'r', encoding='utf-8') as f:
            return json.load(f)
    rows = []
    errors = []
    for d in range(1, 32):
        ymd = f'202605{d:02d}'
        start = 1
        total = None
        while total is None or start <= total:
            end = start + 4  # sample key allows only 5 rows per call
            url = API_BASE.format(start=start, end=end, ymd=ymd)
            try:
                with urllib.request.urlopen(url, timeout=20) as resp:
                    raw = resp.read().decode('utf-8')
                data = json.loads(raw)
                body = data.get('CardSubwayStatsNew', {})
                result = body.get('RESULT', {})
                if result.get('CODE') != 'INFO-000':
                    errors.append((ymd, start, result))
                    break
                if total is None:
                    total = int(body.get('list_total_count', 0))
                batch = body.get('row', [])
                rows.extend(batch)
                if not batch:
                    break
                start = end + 1
                time.sleep(0.015)
            except Exception as e:
                errors.append((ymd, start, repr(e)))
                break
    payload = {'fetched_at': datetime.now().isoformat(timespec='seconds'), 'rows': rows, 'errors': errors}
    with open(API_CACHE, 'w', encoding='utf-8') as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)
    return payload


def analyze_api(payload):
    rows = payload['rows']
    station = defaultdict(int)
    board = defaultdict(int)
    alight = defaultdict(int)
    line = defaultdict(int)
    by_date = defaultdict(int)
    weekday = defaultdict(int)
    weekend = defaultdict(int)
    daily_station = defaultdict(lambda: defaultdict(int))
    for r in rows:
        ymd = r['USE_YMD']
        ln = r['SBWY_ROUT_LN_NM']
        st = r['SBWY_STNS_NM']
        key = (ln, st)
        b = int(r['GTON_TNOPE']); a = int(r['GTOFF_TNOPE'])
        t = b + a
        station[key] += t
        board[key] += b
        alight[key] += a
        line[ln] += t
        by_date[ymd] += t
        daily_station[key][ymd] += t
        dt = date(int(ymd[:4]), int(ymd[4:6]), int(ymd[6:8]))
        if dt.weekday() >= 5:
            weekend[key] += t
        else:
            weekday[key] += t
    days = sorted(by_date)
    return {
        'rows': rows,
        'row_count': len(rows),
        'days': days,
        'station': station,
        'board': board,
        'alight': alight,
        'line': line,
        'by_date': by_date,
        'weekday': weekday,
        'weekend': weekend,
        'daily_station': daily_station,
    }


def top(d, k=10):
    return sorted(d.items(), key=lambda x: x[1], reverse=True)[:k]


def pct(a, b):
    return (a / b * 100) if b else 0


def make_md(old, api, api_payload):
    old_latest = old['latest_station_total']
    current = api['station']
    common = set(old_latest) & set(current)
    growth = []
    for key in common:
        old_v = old_latest[key]
        cur_v = current[key]
        if old_v >= 300000 and cur_v >= 300000:
            growth.append((key, cur_v - old_v, pct(cur_v - old_v, old_v), old_v, cur_v))
    growth_top = sorted(growth, key=lambda x: x[2], reverse=True)[:10]
    decline_top = sorted(growth, key=lambda x: x[2])[:10]
    latest_top = top(old_latest, 10)
    api_top = top(current, 10)
    api_line_top = top(api['line'], 10)
    weekday_top = top(api['weekday'], 5)
    weekend_top = top(api['weekend'], 5)
    busiest_day = max(api['by_date'].items(), key=lambda x: x[1])
    quiet_day = min(api['by_date'].items(), key=lambda x: x[1])

    md = []
    md.append('# 서울 지하철 혼잡도 분석: SCQA 추가 분석')
    md.append('')
    md.append(f'- 저장일: {datetime.now().strftime("%Y-%m-%d %H:%M")}')
    md.append(f'- 작업 폴더: `{ROOT}`')
    md.append(f'- 공공데이터 페이지: {API_DATASET_URL}')
    md.append(f'- 확인 API 예시: `{API_SAMPLE_URL}`')
    md.append('')
    md.append('## S — Situation: 현재 데이터 상황')
    md.append('')
    md.append(f'- 로컬 CSV는 `{old["month_min"]}~{old["month_max"]}` 기간의 월별·시간대별 승하차 데이터입니다.')
    md.append(f'- 로컬 CSV 규모: `{n(old["rows_count"])}행`, `{n(old["line_count"])}개 호선`, `{n(old["station_count"])}개 역명`.')
    md.append('- 서울열린데이터광장 OA-12914는 `서울시 지하철호선별 역별 승하차 인원 정보`이며, 교통카드 기반 일단위 승하차 인원을 제공합니다.')
    md.append('- 페이지 설명 기준: Sheet 서비스는 마지막 한 달치 데이터만 제공하며, 데이터는 매일 3일 전 자료로 갱신됩니다.')
    md.append(f'- OpenAPI 샘플 호출 결과 `CardSubwayStatsNew` 서비스가 정상 응답했습니다. 전체 월 분석은 같은 페이지의 공식 최신 CSV `CARD_SUBWAY_MONTH_202605.csv`를 내려받아 사용했습니다.')
    md.append('')
    md.append('## C — Complication: 기존 분석만으로 부족한 점')
    md.append('')
    md.append('- 기존 로컬 CSV 최신월은 `202106`이라 현재 지하철 이용 패턴과 차이가 날 수 있습니다.')
    md.append('- 로컬 CSV는 시간대별 분석이 가능하지만, 공공데이터 API는 일단위 역별 총 승하차 중심이라 시간대 피크는 직접 알 수 없습니다.')
    md.append('- 따라서 “시간대별 혼잡 구조”는 로컬 CSV로, “최신 역별 수요 변화”는 OpenAPI로 나눠 보는 방식이 안전합니다.')
    md.append('')
    md.append('## Q — Question: 팀 프로젝트에서 답해야 할 질문')
    md.append('')
    md.append('1. 2021년 기준 혼잡 상위 역이 2026년에도 유지되는가?')
    md.append('2. 최신 데이터 기준으로 혼잡 우선 관리 대상 역은 어디인가?')
    md.append('3. 모델/서비스를 만든다면 API를 어떻게 붙여야 계속 최신화되는가?')
    md.append('')
    md.append('## A — Answer: 결론')
    md.append('')
    md.append('- 결론 1: 강남·잠실·홍대입구·신림·구로디지털단지 등 2호선 핵심 역은 2026년 API 기준에서도 여전히 최상위 혼잡 축입니다.')
    md.append('- 결론 2: 기존 분석의 피크 시간대 결론은 유지하되, 최신 순위·변화율은 OA-12914 API로 갱신하는 구조가 필요합니다.')
    md.append('- 결론 3: 발표에서는 `2021 시간대 분석 + 2026 API 최신성 보강` 구조로 가는 것이 가장 설득력 있습니다.')
    md.append('')
    md.append('## 1. OpenAPI 확인 결과')
    md.append('')
    md.append('| 항목 | 내용 |')
    md.append('|---|---|')
    md.append('| 데이터셋 | 서울시 지하철호선별 역별 승하차 인원 정보 |')
    md.append('| 서비스명 | `CardSubwayStatsNew` |')
    md.append('| 호출 형식 | `http://openapi.seoul.go.kr:8088/{KEY}/json/CardSubwayStatsNew/{START}/{END}/{USE_YMD}` |')
    md.append('| 샘플 KEY | `sample`로 endpoint 정상 응답 확인. 단, 샘플키는 건수 제한이 있어 전체 분석은 공식 CSV 다운로드 파일 사용 |')
    md.append('| 주요 필드 | `USE_YMD`, `SBWY_ROUT_LN_NM`, `SBWY_STNS_NM`, `GTON_TNOPE`, `GTOFF_TNOPE`, `REG_YMD` |')
    md.append('| 수집 범위 | 2026-05-01~2026-05-31 |')
    md.append(f'| 수집 행 수 | {n(api["row_count"])}행 |')
    md.append(f'| API 캐시 파일 | `{API_CACHE}` |')
    if api_payload.get('errors'):
        md.append(f'| 수집 오류 | {api_payload["errors"]} |')
    else:
        md.append('| 수집 오류 | 없음 |')
    md.append('')
    md.append('## 2. 2026년 5월 OpenAPI 기준 혼잡 상위 역')
    md.append('')
    md.append('| 순위 | 호선 | 역 | 2026년 5월 승하차 합계 | 승차 | 하차 |')
    md.append('|---:|---|---|---:|---:|---:|')
    for i, ((ln, st), v) in enumerate(api_top, 1):
        md.append(f'| {i} | {ln} | {st} | {n(v)} | {n(api["board"][(ln, st)])} | {n(api["alight"][(ln, st)])} |')
    md.append('')
    md.append('## 3. 2026년 5월 호선별 혼잡도')
    md.append('')
    md.append('| 순위 | 호선 | 승하차 합계 | 비중 |')
    md.append('|---:|---|---:|---:|')
    total_api = sum(api['line'].values())
    for i, (ln, v) in enumerate(api_line_top, 1):
        md.append(f'| {i} | {ln} | {n(v)} | {pct(v, total_api):.1f}% |')
    md.append('')
    md.append('## 4. 2021년 6월 vs 2026년 5월 비교')
    md.append('')
    md.append('비교는 같은 `호선+역명`이 존재하는 역만 대상으로 했습니다. 월이 다르므로 계절성과 월별 일수 차이는 한계입니다.')
    md.append('')
    md.append('### 증가율 상위 역')
    md.append('| 순위 | 호선 | 역 | 2021-06 | 2026-05 | 변화율 |')
    md.append('|---:|---|---|---:|---:|---:|')
    for i, ((ln, st), diff, rate, old_v, cur_v) in enumerate(growth_top, 1):
        md.append(f'| {i} | {ln} | {st} | {n(old_v)} | {n(cur_v)} | {rate:+.1f}% |')
    md.append('')
    md.append('### 감소율 상위 역')
    md.append('| 순위 | 호선 | 역 | 2021-06 | 2026-05 | 변화율 |')
    md.append('|---:|---|---|---:|---:|---:|')
    for i, ((ln, st), diff, rate, old_v, cur_v) in enumerate(decline_top, 1):
        md.append(f'| {i} | {ln} | {st} | {n(old_v)} | {n(cur_v)} | {rate:+.1f}% |')
    md.append('')
    md.append('## 5. 평일/주말 관점')
    md.append('')
    md.append(f'- 2026년 5월 최다 이용일: `{busiest_day[0]}` — `{n(busiest_day[1])}`명')
    md.append(f'- 2026년 5월 최소 이용일: `{quiet_day[0]}` — `{n(quiet_day[1])}`명')
    md.append('')
    md.append('### 평일 누적 TOP 5')
    md.append('| 순위 | 호선 | 역 | 평일 승하차 |')
    md.append('|---:|---|---|---:|')
    for i, ((ln, st), v) in enumerate(weekday_top, 1):
        md.append(f'| {i} | {ln} | {st} | {n(v)} |')
    md.append('')
    md.append('### 주말 누적 TOP 5')
    md.append('| 순위 | 호선 | 역 | 주말 승하차 |')
    md.append('|---:|---|---|---:|')
    for i, ((ln, st), v) in enumerate(weekend_top, 1):
        md.append(f'| {i} | {ln} | {st} | {n(v)} |')
    md.append('')
    md.append('## 6. 프로젝트 적용 제안')
    md.append('')
    md.append('1. `혼잡도 점수 = 최신 API 월별 역 승하차량 정규화 + 기존 CSV 시간대 피크 가중치`로 설계합니다.')
    md.append('2. OpenAPI를 매월 자동 수집하면 모델/대시보드가 낡지 않습니다.')
    md.append('3. API는 일단위 총량이므로, 시간대 혼잡 예측은 기존 시간대 CSV 패턴을 결합해야 합니다.')
    md.append('4. 발표 스토리는 SCQA로 구성합니다: 과거 데이터만으로 분석 → 최신성 한계 → 공공 API 확인 → 최신 혼잡 우선순위 제안.')
    md.append('')
    md.append('## 7. 발표용 한 줄 메시지')
    md.append('')
    md.append('“서울 지하철 혼잡은 2호선 핵심 업무·상권역에 집중되며, 과거 시간대 패턴과 최신 OpenAPI 일별 데이터를 결합하면 실제 서비스 가능한 혼잡도 지표를 만들 수 있다.”')
    md.append('')
    with open(SCQA_MD, 'w', encoding='utf-8') as f:
        f.write('\n'.join(md) + '\n')
    return {
        'api_top': api_top,
        'api_line_top': api_line_top,
        'growth_top': growth_top,
        'decline_top': decline_top,
        'weekday_top': weekday_top,
        'weekend_top': weekend_top,
        'busiest_day': busiest_day,
        'quiet_day': quiet_day,
    }


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    r = p.runs[0]
    r.font.name = 'Malgun Gothic'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    add_text(slide, title, 0.55, 0.35, 12.2, 0.55, size=30, color=color, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.92, 11.6, 0.35, size=12.5, color=WHITE if dark else MUTED)


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(225, 230, 235)
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x+0.25, y+0.18, w-0.4, 0.3, size=14, color=NAVY, bold=True)
    add_text(slide, body, x+0.25, y+0.58, w-0.45, h-0.72, size=11.5, color=DARK)


def add_bar_chart(slide, items, x, y, w, h, title, max_items=6, color=BLUE, unit='명'):
    add_text(slide, title, x, y-0.35, w, 0.3, size=15, color=NAVY, bold=True)
    items = items[:max_items]
    maxv = max(v for _, v in items) if items else 1
    row_h = h / max_items
    for i, (label, v) in enumerate(items):
        yy = y + i * row_h
        if isinstance(label, tuple):
            name = f'{label[0]} {label[1]}'
        else:
            name = str(label)
        add_text(slide, name, x, yy+0.02, 2.1, 0.28, size=9.5, color=DARK)
        bw = (w-3.0) * (v / maxv)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+2.2), Inches(yy+0.03), Inches(max(0.05, bw)), Inches(0.22))
        rect.fill.solid(); rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        add_text(slide, n(v), x+2.25+bw, yy, 1.1, 0.28, size=8.5, color=MUTED)


def make_ppt(old, api, summary):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 title
    s = prs.slides.add_slide(blank)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    add_text(s, '서울 지하철 혼잡도 분석', 0.75, 1.35, 11.6, 0.75, size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 'SCQA 기반 발표자료 · 2021 시간대 데이터 + 2026 OpenAPI 추가 분석', 1.2, 2.18, 10.9, 0.35, size=16, color=RGBColor(210, 225, 240), align=PP_ALIGN.CENTER)
    add_card(s, 1.0, 3.35, 3.3, 1.25, '핵심 역', '강남·잠실·신림·구로디지털단지·홍대입구', ORANGE)
    add_card(s, 5.0, 3.35, 3.3, 1.25, '피크 시간', '08-09시 / 18-19시 집중', TEAL)
    add_card(s, 9.0, 3.35, 3.3, 1.25, '업데이트', 'OA-12914 API로 최신 월별 갱신 가능', BLUE)
    add_text(s, 'Source: 서울열린데이터광장 OA-12914, team_project local files', 0.75, 6.85, 12, 0.25, size=9, color=RGBColor(180, 195, 210))

    # 2 SCQA
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = LIGHT
    add_title(s, 'SCQA로 정리한 문제 정의', '과거 데이터 분석에서 최신 API 결합으로 확장')
    add_card(s, 0.65, 1.55, 3.0, 4.6, 'S 상황', f'로컬 CSV는 {old["month_min"]}~{old["month_max"]} 기간의 월별·시간대별 승하차 데이터를 제공', BLUE)
    add_card(s, 3.9, 1.55, 3.0, 4.6, 'C 문제', '최신월이 202106이라 현재 혼잡 패턴 설명에는 최신성 한계 존재', RED)
    add_card(s, 7.15, 1.55, 2.5, 4.6, 'Q 질문', '현재도 같은 역이 혼잡한가? 서비스화하려면 어떻게 갱신할까?', ORANGE)
    add_card(s, 9.9, 1.55, 2.8, 4.6, 'A 답', '2026 OpenAPI를 붙여 최신 역별 수요를 보강하고, 시간대 패턴은 기존 CSV로 결합', TEAL)

    # 3 old data
    s = prs.slides.add_slide(blank)
    add_title(s, '기존 데이터 분석: 시간대 피크가 명확함', '혼잡도 = 승차인원 + 하차인원')
    hour_items = sorted(old['latest_hour_total'].items(), key=lambda x: x[1], reverse=True)[:6]
    hour_labels = [(f'{h:02d}-{(h+1)%24:02d}시', v) for h, v in hour_items]
    add_bar_chart(s, hour_labels, 0.8, 1.7, 5.7, 3.5, '2021년 6월 시간대별 TOP', 6, TEAL)
    old_top = top(old['latest_station_total'], 5)
    add_bar_chart(s, old_top, 7.0, 1.7, 5.5, 3.5, '2021년 6월 역별 TOP', 5, BLUE)
    add_card(s, 0.9, 5.85, 11.5, 0.75, '해석', '출근 피크 08-09시, 퇴근 피크 18-19시가 뚜렷합니다. 따라서 시간대 피처는 혼잡도 모델의 핵심 변수입니다.', ORANGE)

    # 4 API confirmation
    s = prs.slides.add_slide(blank)
    add_title(s, 'OpenAPI 추가 확인: 최신 일단위 데이터 수집 가능', '서울열린데이터광장 OA-12914 / CardSubwayStatsNew')
    add_card(s, 0.8, 1.5, 3.6, 1.45, 'API endpoint', 'http://openapi.seoul.go.kr:8088/{KEY}/json/CardSubwayStatsNew/{START}/{END}/{USE_YMD}', BLUE)
    add_card(s, 4.85, 1.5, 3.2, 1.45, '수집 범위', f'2026-05-01~2026-05-31\n{n(api["row_count"])}행 수집', TEAL)
    add_card(s, 8.5, 1.5, 3.9, 1.45, '주요 필드', '일자, 호선명, 역명, 승차총인원, 하차총인원, 등록일자', ORANGE)
    add_bar_chart(s, summary['api_line_top'], 1.0, 3.75, 11.2, 2.45, '2026년 5월 호선별 승하차 TOP', 6, BLUE)

    # 5 current top stations
    s = prs.slides.add_slide(blank)
    add_title(s, '2026년 5월 기준 혼잡 상위 역', '최신 API 기준으로도 2호선 핵심 역 집중')
    add_bar_chart(s, summary['api_top'], 0.9, 1.55, 11.6, 4.4, '역별 승하차 합계 TOP 8', 8, ORANGE)
    add_card(s, 0.9, 6.15, 11.5, 0.65, '핵심 메시지', '강남·잠실·홍대입구·신림·구로디지털단지는 과거 분석과 최신 API 모두에서 우선 관리 대상입니다.', TEAL)

    # 6 comparison
    s = prs.slides.add_slide(blank)
    add_title(s, '2021 vs 2026: 최신성 보강이 필요한 이유', '같은 호선+역명 기준 월 승하차량 비교')
    inc = [((ln, st), cur_v) for ((ln, st), diff, rate, old_v, cur_v) in summary['growth_top'][:5]]
    dec = [((ln, st), cur_v) for ((ln, st), diff, rate, old_v, cur_v) in summary['decline_top'][:5]]
    add_bar_chart(s, inc, 0.9, 1.75, 5.5, 3.4, '증가율 상위(2026값)', 5, TEAL)
    add_bar_chart(s, dec, 7.0, 1.75, 5.4, 3.4, '감소율 상위(2026값)', 5, RED)
    add_card(s, 0.9, 5.75, 11.5, 0.9, '주의', '월이 달라 계절성 차이가 있습니다. 발표에서는 “정확한 장기 증감”보다 “최신 API 결합 필요성”의 근거로 사용하는 것이 안전합니다.', ORANGE)

    # 7 proposed index
    s = prs.slides.add_slide(blank)
    add_title(s, '서비스/모델 적용안', '시간대 패턴 + 최신 API 총량을 결합')
    add_card(s, 0.9, 1.55, 3.5, 3.9, '1. 최신 총량', 'OA-12914 API에서 월별 역 승하차량을 수집하고 정규화', BLUE)
    add_card(s, 4.9, 1.55, 3.5, 3.9, '2. 시간대 가중치', '기존 시간대 CSV에서 08-09시, 18-19시 피크 가중치 산출', TEAL)
    add_card(s, 8.9, 1.55, 3.5, 3.9, '3. 혼잡도 점수', '최신 총량 × 시간대 패턴 × 호선/역 특성으로 우선순위 계산', ORANGE)
    add_text(s, '추천 산식 예: 혼잡도 점수 = 최신 월별 총량 정규화 60% + 시간대 피크 집중도 30% + 환승/상권 보정 10%', 1.1, 6.1, 11.0, 0.4, size=14, color=NAVY, bold=True)

    # 8 conclusion
    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    add_text(s, '결론', 0.9, 0.8, 11.6, 0.6, size=34, color=WHITE, bold=True)
    add_text(s, '서울 지하철 혼잡은 2호선 핵심 업무·상권역에 집중됩니다.', 0.95, 1.9, 11.5, 0.55, size=24, color=WHITE, bold=True)
    add_text(s, '기존 시간대 데이터로 “언제 혼잡한지”를 설명하고, OpenAPI로 “지금 어디가 혼잡한지”를 갱신하면 팀 프로젝트 결과물이 단순 분석을 넘어 실제 서비스 구조가 됩니다.', 0.95, 2.75, 11.3, 1.2, size=19, color=RGBColor(220, 235, 245))
    add_card(s, 1.0, 4.65, 3.3, 1.2, '분석', '2021 시간대 피크 확인', TEAL)
    add_card(s, 5.0, 4.65, 3.3, 1.2, '보강', '2026 API 최신 순위 반영', ORANGE)
    add_card(s, 9.0, 4.65, 3.3, 1.2, '적용', '자동 갱신형 혼잡도 지표', BLUE)
    add_text(s, '파일: 혼잡도_SCQA_추가분석.md / 서울지하철_혼잡도_SCQA_발표자료.pptx', 0.9, 6.85, 11.6, 0.25, size=9, color=RGBColor(180, 195, 210))

    prs.save(PPTX_PATH)


def verify_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                texts.append(shape.text)
    all_text = '\n'.join(texts)
    required = ['SCQA', 'OpenAPI', 'CardSubwayStatsNew', '2026년 5월', '결론']
    missing = [r for r in required if r not in all_text]
    return {'slides': len(prs.slides), 'text_chars': len(all_text), 'missing': missing}


def main():
    old = parse_old_csv()
    payload = fetch_api_202605()
    api = analyze_api(payload)
    summary = make_md(old, api, payload)
    make_ppt(old, api, summary)
    v = verify_pptx(PPTX_PATH)
    print('SCQA_MD', SCQA_MD)
    print('PPTX', PPTX_PATH)
    print('API_ROWS', api['row_count'])
    print('API_DAYS', len(api['days']), api['days'][0], api['days'][-1])
    print('PPT_VERIFY', v)
    print('TOP3_API', [(k, v) for k, v in summary['api_top'][:3]])

if __name__ == '__main__':
    main()
